# gradio_app.py
import os
import json
import numpy as np
import faiss
import dashscope
import base64
from http import HTTPStatus
from openai import OpenAI
import gradio as gr
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
from PyPDF2 import PdfReader

# ====================== 全局配置 ======================
DASHSCOPE_API_KEY = os.getenv("DASHSCOPE_API_KEY")
if not DASHSCOPE_API_KEY:
    raise ValueError("请先配置系统环境变量 DASHSCOPE_API_KEY（阿里灵积密钥）")

dashscope.api_key = DASHSCOPE_API_KEY
client = OpenAI(
    api_key=DASHSCOPE_API_KEY,
    base_url="https://dashscope.aliyuncs.com/compatible-mode/v1"
)

MULTIMODAL_EMBEDDING_MODEL = "tongyi-embedding-vision-plus"
INDEX_FILE = "disney_index.faiss"
METADATA_FILE = "disney_metadata.json"
DOCS_DIR = "disney_knowledge_base"
IMG_DIR = os.path.join(DOCS_DIR, "images")
VIDEO_DIR = os.path.join(DOCS_DIR, "videos")

# 检索参数
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
MEDIA_DISTANCE_THRESHOLD = 3.0
IMAGE_KEYWORDS = ["图片", "海报", "照片", "看看", "长什么样", "图"]
VIDEO_KEYWORDS = ["视频", "录像", "影片", "看一下", "播放"]

# 创建文件夹
os.makedirs(DOCS_DIR, exist_ok=True)
os.makedirs(IMG_DIR, exist_ok=True)
os.makedirs(VIDEO_DIR, exist_ok=True)

# ====================== 文档解析器（多格式文档） ======================
def parse_doc(file_path):
    """自动识别文件后缀，解析文本"""
    ext = os.path.splitext(file_path)[1].lower()
    full_text = ""

    if ext in [".doc", ".docx"]:
        doc = DocxDocument(file_path)
        paragraphs = []
        for p in doc.paragraphs:
            if p.text.strip():
                paragraphs.append(p.text.strip())
        # 读取表格
        for table in doc.tables:
            table_lines = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                table_lines.append(row_text)
            paragraphs.append("\n".join(table_lines))
        full_text = "\n".join(paragraphs)

    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            slide_texts.append("\n".join(slide_content))
        full_text = "\n".join(slide_texts)

    elif ext == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            md_raw = f.read()
        full_text = markdown.markdown(md_raw)

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            full_text = f.read()

    elif ext == ".pdf":
        reader = PdfReader(file_path)
        pdf_pages = []
        for page in reader.pages:
            page_txt = page.extract_text()
            if page_txt:
                pdf_pages.append(page_txt)
        full_text = "\n".join(pdf_pages)

    return full_text.strip()

def split_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks

# ====================== 向量生成函数 ======================
def get_text_embedding(text):
    resp = dashscope.MultiModalEmbedding.call(
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{'text': text}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"文本向量失败: {resp.message}")
    return resp.output['embeddings'][0]['embedding']

def get_image_embedding(image_path):
    with open(image_path, "rb") as f:
        base64_image = base64.b64encode(f.read()).decode('utf-8')
    ext = os.path.splitext(image_path)[1].lower().lstrip('.')
    if ext == 'jpg':
        ext = 'jpeg'
    image_data = f"data:image/{ext};base64,{base64_image}"
    resp = dashscope.MultiModalEmbedding.call(
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{'image': image_data}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"图片向量失败: {resp.message}")
    return resp.output['embeddings'][0]['embedding']

def get_video_embedding(video_input):
    """兼容两种输入：网络URL / 本地视频文件路径"""
    if os.path.exists(video_input):
        # 本地视频文件，读取二进制
        with open(video_input, "rb") as f:
            video_bytes = f.read()
        video_base64 = base64.b64encode(video_bytes).decode("utf-8")
        resp = dashscope.MultiModalEmbedding.call(
            model=MULTIMODAL_EMBEDDING_MODEL,
            input=[{'video': video_base64}]
        )
    else:
        # 网络视频URL
        resp = dashscope.MultiModalEmbedding.call(
            model=MULTIMODAL_EMBEDDING_MODEL,
            input=[{'video': video_input}]
        )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"视频向量失败: {resp.message}")
    embeddings = resp.output['embeddings']
    if len(embeddings) > 1:
        vectors = [np.array(e['embedding']) for e in embeddings]
        return np.mean(vectors, axis=0).tolist()
    return embeddings[0]['embedding']

# ====================== 知识库构建（新增本地视频上传处理） ======================
def build_knowledge_base_web(upload_docs, upload_imgs, upload_videos, video_url_input):
    metadata_store = []
    all_vectors = []
    doc_id = 0
    msg_log = []

    # 处理多格式文档
    if upload_docs:
        for doc_file in upload_docs:
            file_name = os.path.basename(doc_file)
            save_path = os.path.join(DOCS_DIR, file_name)
            # 保存本地文件
            with open(save_path, "wb") as f:
                f.write(open(doc_file, "rb").read())
            # 解析全文
            full_text = parse_doc(save_path)
            if not full_text:
                msg_log.append(f"⚠️ {file_name} 未提取到有效文本，跳过")
                continue
            chunks = split_text(full_text)
            msg_log.append(f"✅ 加载文档: {file_name}, 分片数 {len(chunks)}")
            for chunk in chunks:
                metadata = {
                    "id": doc_id,
                    "source": file_name,
                    "type": "text",
                    "content": chunk
                }
                vector = get_text_embedding(chunk)
                all_vectors.append(vector)
                metadata_store.append(metadata)
                doc_id += 1

    # 处理图片
    if upload_imgs:
        for img_file in upload_imgs:
            img_name = os.path.basename(img_file)
            save_path = os.path.join(IMG_DIR, img_name)
            with open(save_path, "wb") as f:
                f.write(open(img_file, "rb").read())
            metadata = {
                "id": doc_id,
                "source": f"图片: {img_name}",
                "type": "image",
                "path": save_path,
                "content": f"[图片素材] {img_name}"
            }
            vector = get_image_embedding(save_path)
            all_vectors.append(vector)
            metadata_store.append(metadata)
            doc_id += 1
        msg_log.append(f"✅ 加载图片 {len(upload_imgs)} 张")

    # 处理本地视频文件
    if upload_videos:
        for vid_file in upload_videos:
            vid_name = os.path.basename(vid_file)
            save_path = os.path.join(VIDEO_DIR, vid_name)
            with open(save_path, "wb") as f:
                f.write(open(vid_file, "rb").read())
            metadata = {
                "id": doc_id,
                "source": f"本地视频: {vid_name}",
                "type": "video",
                "path": save_path,
                "url": "",
                "content": f"[本地视频素材] {vid_name}"
            }
            vector = get_video_embedding(save_path)
            all_vectors.append(vector)
            metadata_store.append(metadata)
            doc_id += 1
        msg_log.append(f"✅ 加载本地视频 {len(upload_videos)} 个")

    # 处理网络视频URL
    if video_url_input and video_url_input.strip():
        metadata = {
            "id": doc_id,
            "source": f"网络视频URL",
            "type": "video",
            "url": video_url_input.strip(),
            "path": "",
            "content": f"[网络视频素材] {video_url_input}"
        }
        vector = get_video_embedding(video_url_input.strip())
        all_vectors.append(vector)
        metadata_store.append(metadata)
        doc_id += 1
        msg_log.append(f"✅ 加载视频URL: {video_url_input}")

    if all_vectors:
        dim = len(all_vectors[0])
        index = faiss.IndexFlatL2(dim)
        index.add(np.array(all_vectors).astype('float32'))
        faiss.write_index(index, INDEX_FILE)
        with open(METADATA_FILE, 'w', encoding='utf-8') as f:
            json.dump(metadata_store, f, ensure_ascii=False, indent=2)
        text_cnt = sum(1 for m in metadata_store if m["type"] == "text")
        img_cnt = sum(1 for m in metadata_store if m["type"] == "image")
        vid_cnt = sum(1 for m in metadata_store if m["type"] == "video")
        msg_log.append(f"\n🎉 知识库构建完成！总条目: {len(metadata_store)} | 文本:{text_cnt} 图片:{img_cnt} 视频:{vid_cnt}")
    else:
        msg_log.append("⚠️ 未检测到任何素材，未生成索引")
    return "\n".join(msg_log)

# ====================== RAG检索问答函数（完全不变） ======================
def load_index():
    index = faiss.read_index(INDEX_FILE)
    with open(METADATA_FILE, 'r', encoding='utf-8') as f:
        metadata = json.load(f)
    return index, metadata

def distance_to_similarity(distance):
    return 1 / (1 + distance)

def detect_media_intent(query):
    query_lower = query.lower()
    want_image = any(kw in query_lower for kw in IMAGE_KEYWORDS)
    want_video = any(kw in query_lower for kw in VIDEO_KEYWORDS)
    return want_image, want_video

def search_with_details(query, index, metadata):
    query_vec = np.array([get_text_embedding(query)]).astype('float32')
    distances, indices = index.search(query_vec, index.ntotal)
    results = []
    for idx, dist in zip(indices[0], distances[0]):
        if idx == -1:
            continue
        m = metadata[idx]
        sim = distance_to_similarity(dist)
        results.append({"idx": idx, "distance": dist, "similarity": sim, "metadata": m})
    return results

def chat_submit(user_msg, chat_history):
    try:
        index, metadata = load_index()
    except Exception as e:
        msg_user = gr.ChatMessage(role="user", content=user_msg)
        msg_ai = gr.ChatMessage(role="assistant", content=f"❌ 知识库加载失败，请先构建索引：{str(e)}")
        chat_history.append(msg_user)
        chat_history.append(msg_ai)
        return chat_history, None, "", "知识库加载失败"

    results = search_with_details(user_msg, index, metadata)
    want_image, want_video = detect_media_intent(user_msg)

    top_texts = [r for r in results if r["metadata"]["type"] == "text"][:3]
    context_str = ""
    for i, r in enumerate(top_texts):
        m = r["metadata"]
        context_str += f"背景知识{i+1}【来源：{m['source']}，相似度：{r['similarity']:.4f}】\n{m['content']}\n\n"

    img_path = None
    if want_image:
        img_candidates = [x for x in results if x["metadata"]["type"] == "image" and x["distance"] < MEDIA_DISTANCE_THRESHOLD]
        if img_candidates:
            img_candidates.sort(key=lambda x: x["distance"])
            img_path = img_candidates[0]["metadata"]["path"]

    video_url = ""
    if want_video:
        vid_candidates = [x for x in results if x["metadata"]["type"] == "video" and x["distance"] < MEDIA_DISTANCE_THRESHOLD]
        if vid_candidates:
            vid_candidates.sort(key=lambda x: x["distance"])
            video_url = vid_candidates[0]["metadata"].get("url", vid_candidates[0]["metadata"].get("path", ""))

    prompt = f"""你是专业问答助手，严格依据背景知识回答，简洁易懂。
【背景知识】
{context_str}
【用户问题】
{user_msg}
"""
    completion = client.chat.completions.create(
        model="qwen-flash",
        messages=[
            {"role": "system", "content": "根据提供的知识库内容精准回答，无相关资料如实说明。"},
            {"role": "user", "content": prompt}
        ]
    )
    answer = completion.choices[0].message.content

    media_tip = ""
    if img_path:
        media_tip += f"\n🖼️ 匹配图片：{os.path.basename(img_path)}"
    if video_url:
        media_tip += f"\n🎬 匹配视频素材：{os.path.basename(video_url)}"
    full_answer = answer + media_tip

    detail_text = "🔍 向量相似度TOP10检索结果：\n"
    for rank, item in enumerate(results[:10]):
        m = item["metadata"]
        detail_text += f"{rank+1}. 【{m['type']}】相似度:{item['similarity']:.4f} 摘要:{m['content'][:45]}...\n"

    chat_history.append(gr.ChatMessage(role="user", content=user_msg))
    chat_history.append(gr.ChatMessage(role="assistant", content=full_answer))
    return chat_history, img_path, video_url, detail_text

# ====================== UI界面（新增本地视频上传组件） ======================
def create_ui():
    with gr.Blocks(title="🎠 多模态RAG智能客服平台") as demo:
        gr.Markdown("""
        # 🎠 多模态RAG智能问答助手
        ### 基于通义多模态向量模型 | 支持文本/图片/视频跨模态检索
        """)
        gr.Markdown("---")

        with gr.Row():
            with gr.Column(scale=3):
                chatbot = gr.Chatbot(
                    label="对话历史",
                    height=480
                )
                user_input = gr.Textbox(
                    label="输入提问",
                    placeholder="例如：查询文档内容、查看图片、查找相关视频",
                    lines=3
                )
                with gr.Row():
                    submit_btn = gr.Button("发送提问", variant="primary")
                    clear_btn = gr.Button("清空对话", variant="secondary")

            with gr.Column(scale=1):
                gr.Markdown("### 📸 匹配素材预览")
                img_preview = gr.Image(label="匹配图片", type="filepath", height=280)
                video_link_box = gr.Textbox(label="匹配视频/文件路径", interactive=False)

        with gr.Accordion("🔍 向量检索详情（调试用）", open=False):
            detail_textbox = gr.Textbox(lines=10, interactive=False, label="相似度排序日志")

        with gr.Accordion("📦 知识库管理 | 上传素材一键建库", open=False):
            gr.Markdown("上传文档/图片/本地视频、填写网络视频URL，支持拖拽上传，构建向量知识库")
            with gr.Row():
                # 文档上传
                upload_docs = gr.File(
                    label="上传文档（doc/docx/ppt/pptx/md/txt/pdf）",
                    file_types=[".doc", ".docx", ".ppt", ".pptx", ".md", ".txt", ".pdf"],
                    file_count="multiple",
                    height=200
                )
                # 图片上传
                upload_imgs = gr.File(
                    label="上传图片",
                    file_types=[".jpg", ".png", ".jpeg", ".webp"],
                    file_count="multiple",
                    height=200
                )
            with gr.Row():
                # 新增本地视频上传
                upload_videos = gr.File(
                    label="上传本地视频（mp4/mov/avi/mkv/flv）",
                    file_types=[".mp4", ".mov", ".avi", ".mkv", ".flv"],
                    file_count="multiple",
                    height=200
                )
                # 保留原有视频URL输入框
                video_url_input = gr.Textbox(label="网络视频URL（公网MP4链接）", placeholder="https://xxx.mp4")
            build_btn = gr.Button("🚀 构建/更新知识库", variant="primary")
            build_log = gr.Textbox(label="构建日志输出", lines=8, interactive=False)

        # 绑定事件（新增upload_videos入参）
        submit_btn.click(
            fn=chat_submit,
            inputs=[user_input, chatbot],
            outputs=[chatbot, img_preview, video_link_box, detail_textbox]
        ).then(lambda: "", outputs=[user_input])

        user_input.submit(
            fn=chat_submit,
            inputs=[user_input, chatbot],
            outputs=[chatbot, img_preview, video_link_box, detail_textbox]
        ).then(lambda: "", outputs=[user_input])

        clear_btn.click(lambda: [], outputs=[chatbot])

        build_btn.click(
            fn=build_knowledge_base_web,
            inputs=[upload_docs, upload_imgs, upload_videos, video_url_input],
            outputs=[build_log]
        )
    return demo

if __name__ == "__main__":
    app = create_ui()
    theme = gr.themes.Soft(
        primary_hue="blue",
        secondary_hue="amber",
        neutral_hue="slate"
    ).set(
        button_primary_background_fill="#165DFF",
        block_radius="12px",
        panel_border_width="1px"
    )
    app.launch(
        server_name="0.0.0.0",
        share=False
    )